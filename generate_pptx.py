from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Setup Presentation (16:9 Widescreen)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# 2. Color Scheme
COLOR_BG = RGBColor(255, 255, 255)
COLOR_TEXT = RGBColor(0, 0, 0)
COLOR_BORDER = RGBColor(0, 0, 0)
COLOR_CONTAINER = RGBColor(245, 247, 250)

# Apply Background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = COLOR_BG

# Helper to add text boxes
def add_label(slide, x, y, w, h, text, size=11, bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = COLOR_TEXT
    p.alignment = align
    return tb

# Helper to add box containers
def add_box(slide, x, y, w, h, text="", bold=False, is_dashed=False, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.color.rgb = COLOR_BORDER
    shape.line.width = Pt(1.5)
    if is_dashed:
        shape.line.dash_style = True
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.bold = bold
        p.font.color.rgb = COLOR_TEXT
        p.alignment = PP_ALIGN.CENTER
    return shape

# --- TITLE ---
add_label(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.5), "Ontology As A Service (OaaS) Architecture", size=20, bold=True, align=PP_ALIGN.LEFT)

# --- LEFT BLOCK: CLIENT CONSUMERS ---
add_box(slide, Inches(0.3), Inches(1.5), Inches(1.5), Inches(4.2))
add_box(slide, Inches(0.5), Inches(1.8), Inches(1.1), Inches(1.5), text="Chatbots", bold=True)
add_box(slide, Inches(0.5), Inches(3.8), Inches(1.1), Inches(1.5), text="Applications", bold=True)

# --- CENTER BLOCK: AMAZON EKS ---
add_box(slide, Inches(2.2), Inches(1.0), Inches(8.2), Inches(5.0))
add_label(slide, Inches(2.4), Inches(1.1), Inches(3.0), Inches(0.4), "Amazon EKS", size=16, bold=True, align=PP_ALIGN.LEFT)

# 1. Ontology Chatbot Backend
add_box(slide, Inches(2.5), Inches(1.6), Inches(3.9), Inches(4.1))
add_label(slide, Inches(2.6), Inches(1.7), Inches(3.5), Inches(0.3), "Ontology Chatbot Backend", size=13, bold=True, align=PP_ALIGN.LEFT)
add_label(slide, Inches(2.6), Inches(2.0), Inches(0.6), Inches(0.2), "pod", size=9)

# Backend components
add_box(slide, Inches(2.7), Inches(2.4), Inches(0.5), Inches(3.1), text="KG API", bold=True) # API column

# Grid coordinates for backend sub-modules
modules_backend = [
    ("Template\nParser", Inches(3.4), Inches(2.4)),
    ("OWL Parser", Inches(4.9), Inches(2.4)),
    ("OWL\nGenerator", Inches(3.4), Inches(3.4)),
    ("Ontology\nManager", Inches(4.9), Inches(3.4)),
    ("Ontology\nPublisher", Inches(3.4), Inches(4.4)),
    ("Ontology\nChatbot", Inches(4.9), Inches(4.4))
]
for text, x, y in modules_backend:
    add_box(slide, x, y, Inches(1.3), Inches(0.8), text=text)

# 2. Ontology Chatbot Frontend
add_box(slide, Inches(6.7), Inches(1.6), Inches(3.4), Inches(4.1))
add_label(slide, Inches(6.8), Inches(1.7), Inches(3.0), Inches(0.3), "Ontology Chatbot Frontend", size=13, bold=True, align=PP_ALIGN.LEFT)
add_label(slide, Inches(6.8), Inches(2.0), Inches(0.6), Inches(0.2), "pod", size=9)

modules_frontend = [
    ("Persona\nmanager", Inches(6.9), Inches(2.5)),
    ("Ontology\nManager", Inches(8.4), Inches(2.5)),
    ("Ontology\nGraph Viewer", Inches(6.9), Inches(3.8)),
    ("Ontology\nChatbot", Inches(8.4), Inches(3.8))
]
for text, x, y in modules_frontend:
    add_box(slide, x, y, Inches(1.3), Inches(1.0), text=text)


# --- RIGHT BLOCK: ROLES / USERS ---
roles = ["Domain Owner", "Developers", "Users"]
for i, role in enumerate(roles):
    y_offset = Inches(1.5 + (i * 1.5))
    add_label(slide, Inches(10.8), y_offset + Inches(0.4), Inches(2.0), Inches(0.4), role, size=13, bold=True)

# --- BOTTOM BLOCK: PERSISTENCE & DATA PIPELINES ---
# Amazon Neptune
add_box(slide, Inches(2.7), Inches(6.6), Inches(1.5), Inches(0.6), text="Amazon\nNeptune")

# oaas_ontology* Container
add_box(slide, Inches(4.6), Inches(6.4), Inches(3.2), Inches(0.9), fill_color=COLOR_CONTAINER)
add_label(slide, Inches(4.7), Inches(6.5), Inches(1.5), Inches(0.3), "oaas_ontology*", size=11)
add_label(slide, Inches(6.3), Inches(6.5), Inches(1.4), Inches(0.3), "Ontology agent", size=11)

# Databricks
add_box(slide, Inches(8.2), Inches(6.6), Inches(1.5), Inches(0.6), text="Databricks")

# Save Presentation
prs.save("Ontology_Architecture.pptx")
print("Successfully generated editable PowerPoint presentation framework.")
